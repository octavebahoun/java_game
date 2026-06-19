from pptx import Presentation
import sys

path = sys.argv[1] if len(sys.argv) > 1 else 'Jeu_Mots_Caches_Pro.pptx'
prs = Presentation(path)
print(f'Slides: {len(prs.slides)}')
for i, slide in enumerate(prs.slides, start=1):
    title = ''
    # Try canonical title placeholder
    try:
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text.strip()
    except Exception:
        title = ''
    # Fallback: first non-empty text frame in shapes
    if not title:
        for shape in slide.shapes:
            try:
                if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                    txt = shape.text_frame.text.strip()
                    if txt:
                        title = txt.split('\n')[0].strip()
                        break
            except Exception:
                continue
    notes = ''
    try:
        notes = slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        notes = ''
    print(f'--- Slide {i} ---')
    print('Title:', title)
    print('Notes:', notes)
