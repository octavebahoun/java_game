from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# Helper to add slide with title and bullet points and speaker notes
def add_bullet_slide(title, bullets, notes=None):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = b
        else:
            p = tf.add_paragraph()
            p.text = b
            p.level = 1
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide

# Title slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Jeu de Mots Cachés"
slide.placeholders[1].text = "Présentation des fonctions principales"
slide.notes_slide.notes_text_frame.text = "Bonjour, je présente le projet 'Jeu de Mots Cachés'. Aujourd'hui, nous allons suivre le flux logique du jeu en explorant les fonctions principales qui le font fonctionner."

# Contexte
add_bullet_slide(
    "Contexte et objectif",
    [
        "Jeu éducatif en Java : trouver des mots cachés dans une grille 10x10",
        "Deux modes : console (CLI) et interface graphique (Swing)",
        "Architecture MVC : séparation model / logique / interface"
    ],
    notes="Le projet suit une architecture claire avec packages distincts : model, moteur, view, exception."
)

# Flux global
add_bullet_slide(
    "Flux global du jeu",
    [
        "1. Initialiser une partie → preparerPartie() place les mots",
        "2. Afficher la grille → getGrille().toString()",
        "3. Lire les saisies du joueur → lireEntier()",
        "4. Vérifier le mot → verifier() [fonction PRINCIPALE]",
        "5. Mettre à jour l'état → getDictionnaire(), getScore()"
    ],
    notes="Voici le flux du jeu étape par étape. Chaque étape correspond à une ou plusieurs fonctions que nous allons explorer en détail."
)

# Étape 1 : Initialisation
add_bullet_slide(
    "Étape 1 : Initialiser la partie",
    [
        "new MoteurJeu() crée une instance du moteur",
        "preparerPartie() place 10 mots dans la grille",
        "placerMot(String, ligne, col, dLigne, dCol) : place un mot et vérifie les collisions",
        "remplirCasesVides() complète avec des lettres aléatoires"
    ],
    notes="Au démarrage, la grille est créée vide. On place d'abord les mots de la partie (CHIEN, PYTHON, etc), puis on remplit le reste aléatoirement. Chaque mot est aussi enregistré dans le dictionnaire."
)

# Étape 2 : Afficher
add_bullet_slide(
    "Étape 2 : Afficher la grille",
    [
        "getGrille() récupère l'objet Grille du moteur",
        "toString() affiche la grille avec numéros de lignes/colonnes",
        "getCase(ligne, colonne) retourne la lettre à une position",
        "Utile : le joueur doit lire les coordonnées sur la grille"
    ],
    notes="L'affichage est simple : tableau 10x10 avec en-têtes. Les numéros aident le joueur à repérer les mots visuellement et à noter les coordonnées de début et de fin."
)

# Étape 3 : Lire saisies
add_bullet_slide(
    "Étape 3 : Récupérer les saisies du joueur",
    [
        "lireEntier(Scanner, String) : lit et parse un entier au clavier",
        "Récupère 4 valeurs : ligne début, colonne début, ligne fin, colonne fin",
        "Crée une Coordonnee(ligne, colonne) pour chaque point",
        "Gère les erreurs : si l'user tape des lettres → NumberFormatException"
    ],
    notes="Quatre appels à lireEntier() pour récupérer les 4 coordonnées. La fonction parse la saisie et la convertit en entier, sinon elle lève une exception qui sera attrapée au niveau supérieur."
)

# Étape 4 : Vérifier (PRINCIPALE)
add_bullet_slide(
    "Étape 4a : Vérifier les coordonnées",
    [
        "verifierDansGrille(Coordonnee) : valide que (l,c) est dans [0,9]",
        "Lance CoordonneesInvalidesException si hors limites",
        "Exemple : si ligne = 15 → erreur (max = 9)",
        "C'est une barrière de sécurité avant de lire la grille"
    ],
    notes="Cette vérification évite les IndexOutOfBoundsException et donne au joueur un message clair. Elle est appelée deux fois : une pour debut, une pour fin."
)

add_bullet_slide(
    "Étape 4b : Vérifier l'alignement",
    [
        "Calcule distLigne = |fin.ligne - debut.ligne|",
        "Calcule distColonne = |fin.colonne - debut.colonne|",
        "Horizontal : distLigne == 0 && distColonne != 0",
        "Vertical : distColonne == 0 && distLigne != 0",
        "Diagonal : distLigne == distColonne && distLigne != 0",
        "Sinon → AlignementInvalideException (ex: déplacement en L)"
    ],
    notes="L'algorithme rejet tout ce qui n'est pas une ligne droite. Par exemple, aller de (0,0) à (1,2) est un L invalide. Seul H/V/D sont acceptés."
)

add_bullet_slide(
    "Étape 4c : Lire le mot dans la grille",
    [
        "Calcule le pas : pasLigne, pasColonne = -1, 0 ou +1",
        "Boucle de debut jusqu'à fin, appelle getCase(l,c)",
        "Accumule les lettres dans motLu",
        "Compare motLu avec le mot saisi (insensible à la casse)",
        "Vérifie aussi la lecture en INVERSE (mots lus à l'envers)"
    ],
    notes="On lit lettre par lettre en suivant la direction calculée. Exemple : de (0,0) à (0,4) horizontal lit A, B, C, D, E. On compare aussi l'inverse au cas où le mot serait marqué de droite à gauche."
)

add_bullet_slide(
    "Étape 4d : Vérifier dans le dictionnaire",
    [
        "dictionnaire.contient(motCherche) : le mot est-il à trouver ?",
        "dictionnaire.estDejaTrouve(mot) : a-t-il déjà été trouvé ?",
        "Si trouvé ET dans dictionnaire → dictionnaire.supprimer(motCherche)",
        "Met à jour le score : score += motCherche.length() * 10"
    ],
    notes="Le dictionnaire maintient deux listes : la liste des mots restants et la liste complète. Quand un mot est trouvé, on le supprime de la liste restante."
)

# Étape 5 : Fin et boucle
add_bullet_slide(
    "Étape 5 : Vérifier la fin de partie",
    [
        "partieGagnee() : dictionnaire.estVide() ?",
        "Si oui → affiche le score final et fin du jeu",
        "Sinon → boucle retour à l'étape 2 (afficher grille)",
        "Affiche aussi les mots trouvés (barrés) et restants"
    ],
    notes="Le jeu boucle tant qu'il reste des mots. À chaque itération, le joueur voit la grille avec les mots trouvés colorisés, le score et la liste des mots restants."
)

# Gestion des erreurs
add_bullet_slide(
    "Gestion des erreurs avec try-catch",
    [
        "CoordonneesInvalidesException : coordonnées hors grille",
        "AlignementInvalideException : pas une ligne droite",
        "NumberFormatException : saisie non numérique",
        "Chaque exception affiche un message clair au joueur",
        "Le jeu ne crash jamais : boucle continue après erreur"
    ],
    notes="Les trois exceptions couvertes permettent une expérience robuste. Le joueur peut corriger sa saisie sans perdre la partie."
)

# Architecture logicielle
add_bullet_slide(
    "Packages et responsabilités",
    [
        "model : Grille, Dictionnaire, Coordonnee (données)",
        "moteur : MoteurJeu, AppConsole (logique du jeu)",
        "view : FenetrePrincipale (interface graphique Swing)",
        "exception : CoordonneesInvalidesException, AlignementInvalideException"
    ],
    notes="Cette séparation permet de tester la logique sans interface, et de changer l'interface sans modifier le moteur."
)

# Interface graphique (bonus)
add_bullet_slide(
    "Interface graphique (FenetrePrincipale)",
    [
        "Grille visuelle : GridLayout 10x10 avec cases colorées",
        "Champs de texte pour mot et 4 coordonnées",
        "Boutons : Valider (appelle onValider()) et Rejouer",
        "Mots trouvés sont barrés et colorisés en vert",
        "Popups d'erreur pour exceptions et feedback utilisateur"
    ],
    notes="La GUI utilise la même logique MoteurJeu et vérifier(). Elle enrichit l'UX avec des visuels et des retours immédiats."
)

# Conclusion
add_bullet_slide(
    "Conclusion",
    [
        "Les fonctions principales : preparerPartie(), getCase(), lireEntier(), verifier()",
        "verifier() est le cœur : validation, calcul alignement, lecture, vérification",
        "Architecture claire et extensible (ajouter modes, listes externes, etc.)",
        "Robustesse : try-catch sur toutes les saisies utilisateur"
    ],
    notes="Merci ! Avez-vous des questions sur les fonctions ou le flux du jeu ?"
)

# Save
prs.save('presentation.pptx')
print('presentation.pptx créée avec succès.')
